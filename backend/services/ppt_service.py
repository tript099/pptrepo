from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os
import uuid
from typing import Dict, Any, List, Optional
from datetime import datetime
import tempfile
import shutil

class PPTService:
    def __init__(self):
        self.output_dir = "outputs"
        os.makedirs(self.output_dir, exist_ok=True)
        self.template_assets = {}  # Store extracted template assets

    def analyze_template(self, template_path: str):
        """
        Comprehensive template analysis to extract all design elements
        """
        print("🔍 Analyzing template comprehensively...")
        
        try:
            template_prs = Presentation(template_path)
            self.template_assets = {
                'logos': [],
                'backgrounds': [],
                'slide_masters': [],
                'color_schemes': [],
                'fonts': [],
                'layouts': []
            }
            
            # 1. Extract slide master information
            slide_master = template_prs.slide_master
            print(f"📋 Slide Master: {slide_master}")
            
            # 2. Extract images from slide master first (most important)
            print("🖼️  Extracting images from slide master...")
            try:
                for shape in slide_master.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Extract image data
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Save the image temporarily
                            temp_image_path = os.path.join(tempfile.gettempdir(), f"master_image_{len(self.template_assets['logos'])}.{image.ext}")
                            with open(temp_image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            # Store image info
                            image_info = {
                                'path': temp_image_path,
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'slide_index': 'master',
                                'filename': f"master_logo_{len(self.template_assets['logos'])}.{image.ext}",
                                'is_master': True
                            }
                            self.template_assets['logos'].append(image_info)
                            print(f"   📸 Found MASTER image: {image_info['filename']} at position ({shape.left}, {shape.top})")
                            
                        except Exception as e:
                            print(f"   ❌ Could not extract master image: {e}")
            except Exception as e:
                print(f"   ❌ Could not access slide master shapes: {e}")
            
            # 3. Extract images from slide layouts
            print("🖼️  Extracting images from slide layouts...")
            for layout_idx, layout in enumerate(template_prs.slide_layouts):
                try:
                    for shape in layout.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                # Extract image data
                                image = shape.image
                                image_bytes = image.blob
                                
                                # Save the image temporarily
                                temp_image_path = os.path.join(tempfile.gettempdir(), f"layout_{layout_idx}_image_{len(self.template_assets['logos'])}.{image.ext}")
                                with open(temp_image_path, 'wb') as f:
                                    f.write(image_bytes)
                                
                                # Store image info
                                image_info = {
                                    'path': temp_image_path,
                                    'left': shape.left,
                                    'top': shape.top,
                                    'width': shape.width,
                                    'height': shape.height,
                                    'slide_index': f'layout_{layout_idx}',
                                    'filename': f"layout_{layout_idx}_logo_{len(self.template_assets['logos'])}.{image.ext}",
                                    'is_layout': True
                                }
                                self.template_assets['logos'].append(image_info)
                                print(f"   📸 Found LAYOUT image: {image_info['filename']} at position ({shape.left}, {shape.top})")
                                
                            except Exception as e:
                                print(f"   ❌ Could not extract layout image: {e}")
                except Exception as e:
                    print(f"   ❌ Could not access layout {layout_idx} shapes: {e}")
            
            # 4. Extract all images/logos from template slides as fallback
            print("🖼️  Extracting images from template slides...")
            for slide_idx, slide in enumerate(template_prs.slides):
                print(f"   🔍 Checking slide {slide_idx} for images...")
                shape_count = 0
                for shape in slide.shapes:
                    shape_count += 1
                    print(f"      Shape {shape_count}: Type={shape.shape_type}")
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Extract image data
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Save the image temporarily
                            temp_image_path = os.path.join(tempfile.gettempdir(), f"slide_{slide_idx}_image_{len(self.template_assets['logos'])}.{image.ext}")
                            with open(temp_image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            # Store image info
                            image_info = {
                                'path': temp_image_path,
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height,
                                'slide_index': slide_idx,
                                'filename': f"slide_{slide_idx}_logo_{len(self.template_assets['logos'])}.{image.ext}",
                                'is_slide': True
                            }
                            self.template_assets['logos'].append(image_info)
                            print(f"   📸 Found SLIDE image: {image_info['filename']} at position ({shape.left}, {shape.top})")
                            
                        except Exception as e:
                            print(f"   ❌ Could not extract slide image: {e}")
                            import traceback
                            traceback.print_exc()
                print(f"   📊 Slide {slide_idx} has {shape_count} total shapes")
            
            # 5. Extract background information
            print("🎨 Extracting background information...")
            for layout_idx, layout in enumerate(template_prs.slide_layouts):
                try:
                    background_info = {
                        'layout_index': layout_idx,
                        'layout_name': layout.name,
                        'background': getattr(layout, 'background', None)
                    }
                    self.template_assets['backgrounds'].append(background_info)
                    print(f"   🎭 Layout {layout_idx}: {layout.name}")
                except Exception as e:
                    print(f"   ❌ Could not extract background for layout {layout_idx}: {e}")
            
            print(f"✅ Template analysis complete:")
            print(f"   📸 Images/Logos found: {len(self.template_assets['logos'])}")
            print(f"   🎭 Layouts found: {len(self.template_assets['backgrounds'])}")
            
            if len(self.template_assets['logos']) == 0:
                print("🚨 WARNING: No images/logos found in template!")
                print("   This might mean:")
                print("   - Template has no embedded images")
                print("   - Images are in slide master (need different extraction)")
                print("   - Images are part of background (need different approach)")
            
            return self.template_assets
            
        except Exception as e:
            print(f"❌ Template analysis failed: {e}")
            import traceback
            traceback.print_exc()
            return {}

    def create_presentation_with_full_template(self, slide_data: Dict[str, Any], template_path: str = None, logo_path: str = None, logo_position: str = "top-right", logo_size: str = "medium", output_dir: str = None) -> str:
        """
        Create presentation with comprehensive template inheritance
        """
        print("🚀 Creating presentation with full template inheritance...")
        
        # Step 1: Analyze template if provided
        if template_path and os.path.exists(template_path):
            template_assets = self.analyze_template(template_path)
            
            # Check if this is an original template from extracted presentation
            is_original_template = slide_data.get("meta", {}).get("has_template", False)
            
            if is_original_template:
                print("🔄 Using original presentation as template for editing...")
                # Use the original presentation directly and replace content
                prs = Presentation(template_path)
                template_slide_count = len(prs.slides)
                need_slide_cleanup = True
                
                # We'll replace existing slides with edited content
                print(f"📋 Original presentation has {template_slide_count} slides")
                
            else:
                # Create a clean presentation but preserve template structure
                print("📋 Creating clean presentation with template structure...")
                template_prs = Presentation(template_path)
                
                # Create new presentation with template structure
                prs = Presentation()
                
                # Copy slide master and layouts from template
                try:
                    print("🔄 Copying slide master and layouts...")
                    prs.slide_master._element = template_prs.slide_master._element
                    prs.slide_layouts._sldLayoutLst = template_prs.slide_layouts._sldLayoutLst
                    print("✅ Successfully copied template structure")
                except Exception as e:
                    print(f"❌ Could not copy template structure: {e}")
                    # Fallback: use template directly and remove slides later
                    prs = template_prs
                    template_slide_count = len(prs.slides)
                    need_slide_cleanup = True
                else:
                    need_slide_cleanup = False
                    template_slide_count = 0
                
        else:
            print("📄 Using default blank presentation")
            prs = Presentation()
            template_assets = {}
            need_slide_cleanup = False
            template_slide_count = 0
            is_original_template = False
        
        print(f"📊 Starting to create {len(slide_data['slides'])} slides...")
        
        # Step 2: Create/Update slides
        if is_original_template and len(prs.slides) == len(slide_data["slides"]):
            print("🔄 Updating existing slides with edited content...")
            
            # Update existing slides instead of creating new ones
            for i, slide_info in enumerate(slide_data["slides"]):
                print(f"📝 Updating slide {i+1}: {slide_info.get('title', 'Untitled')}")
                
                try:
                    existing_slide = prs.slides[i]
                    
                    # Clear existing content but keep layout and background
                    self._update_existing_slide_content(existing_slide, slide_info)
                    
                    # Add uploaded logo if provided
                    if logo_path and os.path.exists(logo_path):
                        self._add_logo_to_slide(existing_slide, logo_path, logo_position, logo_size)
                        
                except Exception as e:
                    print(f"❌ Error updating slide {i+1}: {e}")
                    import traceback
                    traceback.print_exc()
        else:
            print("📄 Creating new slides...")
            
            # Create new AI-generated slides
            for i, slide_info in enumerate(slide_data["slides"]):
                print(f"📄 Creating slide {i+1}: {slide_info.get('title', 'Untitled')}")
                
                try:
                    if slide_info["layout"] == "title":
                        slide = self._create_title_slide_with_template(prs, slide_info, template_assets)
                    elif slide_info["layout"] == "bullets":
                        slide = self._create_bullet_slide_with_template(prs, slide_info, template_assets)
                    elif slide_info["layout"] == "table":
                        slide = self._create_table_slide_with_template(prs, slide_info, template_assets)
                    elif slide_info["layout"].startswith("chart"):
                        slide = self._create_chart_slide_with_template(prs, slide_info, template_assets)
                    else:
                        print(f"⚠️  Unknown layout: {slide_info['layout']}, using bullet layout")
                        slide = self._create_bullet_slide_with_template(prs, slide_info, template_assets)
                    
                    # Add template assets (logos, etc.) to each slide
                    self._apply_template_assets_to_slide(slide, template_assets)
                    
                    # Add any text boxes from extracted data
                    self._add_text_boxes_to_slide(slide, slide_info)
                    
                    # Add uploaded logo if provided
                    if logo_path and os.path.exists(logo_path):
                        self._add_logo_to_slide(slide, logo_path, logo_position, logo_size)
                        
                except Exception as e:
                    print(f"❌ Error creating slide {i+1}: {e}")
                    import traceback
                    traceback.print_exc()
        
        # Step 3: Clean up template slides if needed
        if need_slide_cleanup and template_slide_count > 0:
            print(f"🧹 Cleaning up {template_slide_count} template slides...")
            ai_slide_count = len(slide_data["slides"])
            
            # Remove template slides from the beginning
            slides_removed = 0
            for i in range(template_slide_count):
                if len(prs.slides) > ai_slide_count:
                    try:
                        slide_elem = prs.slides._sldIdLst[0]
                        prs.slides._sldIdLst.remove(slide_elem)
                        slides_removed += 1
                    except Exception as e:
                        print(f"❌ Error removing template slide {i+1}: {e}")
                        break
                else:
                    break
            
            print(f"✅ Removed {slides_removed} template slides")
        
        print(f"🎯 Final presentation: {len(prs.slides)} slides")
        
        # Step 4: Save presentation
        deck_title = slide_data["meta"]["deck_title"].replace(" ", "_").replace("/", "_")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{deck_title}_{timestamp}.pptx"
        
        output_path = os.path.join(output_dir or self.output_dir, filename)
        prs.save(output_path)
        
        print(f"💾 Presentation saved: {output_path}")
        return output_path

    def _get_slide_layout(self, prs, layout_type: str):
        """
        Intelligently find the best slide layout based on type and available layouts
        """
        available_layouts = len(prs.slide_layouts)
        print(f"Template has {available_layouts} slide layouts available")
        
        # Log all available layouts for debugging
        layout_names = []
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = layout.name.lower()
            layout_names.append(layout_name)
            print(f"Layout {i}: {layout.name}")
        
        try:
            if layout_type == 'title':
                # Look for title-specific layouts first
                for i, name in enumerate(layout_names):
                    if 'title' in name and 'slide' in name:
                        print(f"Found title layout at index {i}: {prs.slide_layouts[i].name}")
                        return prs.slide_layouts[i]
                # Fallback to first layout
                return prs.slide_layouts[0]
                    
            elif layout_type == 'content':
                # Look for content, bullet, or text layouts
                for i, name in enumerate(layout_names):
                    if any(keyword in name for keyword in ['content', 'bullet', 'text', 'list']):
                        print(f"Found content layout at index {i}: {prs.slide_layouts[i].name}")
                        return prs.slide_layouts[i]
                # Fallback to second layout if available, then first
                if available_layouts > 1:
                    return prs.slide_layouts[1]
                return prs.slide_layouts[0]
                    
            elif layout_type == 'blank':
                # Look for blank layouts first
                for i, name in enumerate(layout_names):
                    if 'blank' in name:
                        print(f"Found blank layout at index {i}: {prs.slide_layouts[i].name}")
                        return prs.slide_layouts[i]
                # Look for content layouts as alternative
                for i, name in enumerate(layout_names):
                    if any(keyword in name for keyword in ['content', 'text']):
                        print(f"Using content layout for blank at index {i}: {prs.slide_layouts[i].name}")
                        return prs.slide_layouts[i]
                # Fallback to last available layout, then first
                if available_layouts > 2:
                    return prs.slide_layouts[-1]
                elif available_layouts > 1:
                    return prs.slide_layouts[1]
                return prs.slide_layouts[0]
                    
        except IndexError:
            pass
        
        # Ultimate fallback - use the first available layout
        if available_layouts > 0:
            print(f"Using fallback layout 0 for {layout_type}")
            return prs.slide_layouts[0]
        else:
            raise Exception("No slide layouts available in template")

    def _add_logo_to_slide(self, slide, logo_path: str, position: str, size: str):
        """Add logo to a slide at specified position and size"""
        try:
            # Size mapping
            size_map = {
                "small": 1.0,
                "medium": 1.5,
                "large": 2.0
            }
            logo_size = size_map.get(size, 1.5)
            
            # Position mapping (in inches from slide edges)
            slide_width = 10  # Standard slide width
            slide_height = 7.5  # Standard slide height
            margin = 0.3
            
            position_map = {
                "top-left": (margin, margin),
                "top-right": (slide_width - logo_size - margin, margin),
                "bottom-left": (margin, slide_height - logo_size - margin),
                "bottom-right": (slide_width - logo_size - margin, slide_height - logo_size - margin),
                "center": ((slide_width - logo_size) / 2, (slide_height - logo_size) / 2)
            }
            
            left, top = position_map.get(position, position_map["top-right"])
            
            # Add logo to slide
            slide.shapes.add_picture(
                logo_path,
                Inches(left),
                Inches(top),
                width=Inches(logo_size),
                height=Inches(logo_size)
            )
            print(f"Added logo to slide at position {position} with size {size}")
            
        except Exception as e:
            print(f"Warning: Could not add logo to slide: {e}")
            # Continue without logo rather than failing
    
    def _apply_corporate_template(self, prs):
        """Apply corporate styling to the presentation"""
        # This would typically load a custom template
        # For now, we'll apply some basic corporate styling
        pass
    
    def _create_title_slide(self, prs, slide_info: Dict[str, Any]):
        """Create a title slide with proper template placeholder detection"""
        slide_layout = self._get_slide_layout(prs, 'title')
        slide = prs.slides.add_slide(slide_layout)
        
        print(f"Title slide layout: {slide_layout.name}")
        print(f"Available placeholders: {len(slide.placeholders)}")
        
        # List all placeholders for debugging
        for i, placeholder in enumerate(slide.placeholders):
            placeholder_type = placeholder.placeholder_format.type
            placeholder_name = getattr(placeholder, 'name', 'Unknown')
            print(f"  Placeholder {i}: Type={placeholder_type}, Name={placeholder_name}")
        
        # Set title using the most appropriate method
        title_shape = None
        try:
            # Method 1: Use the title shape if available
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide.shapes.title.text = slide_info["title"]
                title_shape = slide.shapes.title
                print("Used slide.shapes.title for title")
            else:
                raise AttributeError("No title shape available")
        except AttributeError:
            # Method 2: Find title placeholder by type
            title_set = False
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:  # Title placeholder
                    try:
                        placeholder.text = slide_info["title"]
                        title_shape = placeholder
                        print("Used title placeholder (type 1) for title")
                        title_set = True
                        break
                    except Exception as e:
                        print(f"Could not use title placeholder: {e}")
            
            # Method 3: Use first placeholder as fallback
            if not title_set and len(slide.placeholders) > 0:
                try:
                    slide.placeholders[0].text = slide_info["title"]
                    title_shape = slide.placeholders[0]
                    print("Used first placeholder for title")
                except Exception as e:
                    print(f"Could not use first placeholder: {e}")
                    # Method 4: Create text box as last resort
                    try:
                        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
                        title_box.text_frame.text = slide_info["title"]
                        title_shape = title_box
                        print("Created text box for title")
                    except Exception as e2:
                        print(f"Could not create title text box: {e2}")
        
        # Set subtitle if present
        if "subtitle" in slide_info and slide_info["subtitle"]:
            subtitle_set = False
            
            # Method 1: Find subtitle placeholder by type
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 3:  # Subtitle placeholder
                    try:
                        placeholder.text = slide_info["subtitle"]
                        print("Used subtitle placeholder (type 3) for subtitle")
                        subtitle_set = True
                        break
                    except Exception as e:
                        print(f"Could not use subtitle placeholder: {e}")
            
            # Method 2: Use second placeholder as fallback
            if not subtitle_set and len(slide.placeholders) > 1:
                try:
                    slide.placeholders[1].text = slide_info["subtitle"]
                    print("Used second placeholder for subtitle")
                    subtitle_set = True
                except Exception as e:
                    print(f"Could not use second placeholder: {e}")
            
            # Method 3: Create text box for subtitle
            if not subtitle_set:
                try:
                    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
                    subtitle_box.text_frame.text = slide_info["subtitle"]
                    print("Created text box for subtitle")
                except Exception as e:
                    print(f"Could not create subtitle text box: {e}")
        
        # Apply styling only if we have a title shape
        if title_shape:
            self._style_title_text(title_shape)
        
        return slide
    
    def _create_bullet_slide(self, prs, slide_info: Dict[str, Any]):
        """Create a bullet point slide"""
        # Find the best content layout
        slide_layout = self._get_slide_layout(prs, 'content')
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        print(f"Bullet slide layout: {slide_layout.name}")
        print(f"Available placeholders: {len(slide.placeholders)}")
        
        try:
            title = slide.shapes.title
            title.text = slide_info["title"]
            print("Used slide.shapes.title for bullet slide title")
            self._style_title_text(title)
        except Exception as e:
            print(f"Could not use slide.shapes.title: {e}")
            # Try first placeholder
            if len(slide.placeholders) > 0:
                try:
                    slide.placeholders[0].text = slide_info["title"]
                    title = slide.placeholders[0]
                    print("Used first placeholder for bullet slide title")
                    self._style_title_text(title)
                except Exception as e2:
                    print(f"Could not use first placeholder: {e2}")
        
        # Set bullet points
        try:
            # List all placeholders for debugging
            for i, placeholder in enumerate(slide.placeholders):
                placeholder_type = placeholder.placeholder_format.type
                placeholder_name = getattr(placeholder, 'name', 'Unknown')
                print(f"  Placeholder {i}: Type={placeholder_type}, Name={placeholder_name}")
            
            # Try to find a content placeholder
            content_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 2:  # Content placeholder
                    content_placeholder = placeholder
                    print(f"Found content placeholder (type 2)")
                    break
            
            if content_placeholder is None and len(slide.placeholders) > 1:
                # Fallback to second placeholder if available
                content_placeholder = slide.placeholders[1]
                print(f"Using second placeholder as content placeholder")
            
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for i, bullet in enumerate(slide_info["bullets"]):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet
                    p.level = 0
                    self._style_bullet_text(p)
            else:
                raise Exception("No suitable placeholder found")
                
        except Exception as e:
            print(f"Warning: Could not add bullet points to placeholder: {e}")
            # Create a text box as fallback
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = text_box.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(slide_info["bullets"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = f"• {bullet}"
                self._style_bullet_text(p)
        
        return slide
    
    def _create_table_slide(self, prs, slide_info: Dict[str, Any]):
        """Create a table slide"""
        # Find the best blank or content layout for tables
        slide_layout = self._get_slide_layout(prs, 'blank')
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = slide_info["title"]
        self._style_title_text(title_shape)
        
        # Calculate table dimensions
        rows = len(slide_info["rows"]) + 1  # +1 for header
        cols = len(slide_info["columns"])
        
        # Add table
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4)
        )
        table = table_shape.table
        
        # Set column headers
        for col_idx, col_name in enumerate(slide_info["columns"]):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            self._style_table_header(cell)
        
        # Set data rows
        for row_idx, row_data in enumerate(slide_info["rows"]):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                self._style_table_cell(cell)
        
        return slide
    
    def _create_chart_slide(self, prs, slide_info: Dict[str, Any], chart_type: str):
        """Create a chart slide"""
        # Find the best blank or content layout for charts
        slide_layout = self._get_slide_layout(prs, 'blank')
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = slide_info["title"]
        self._style_title_text(title_shape)
        
        # Prepare chart data
        chart_data = CategoryChartData()
        chart_data.categories = slide_info["categories"]
        
        for series_name, values in slide_info["series"].items():
            chart_data.add_series(series_name, values)
        
        # Map chart types
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Add chart
        chart_shape = slide.shapes.add_chart(
            xl_chart_type, Inches(1), Inches(1.5), Inches(8), Inches(4), chart_data
        )
        
        self._style_chart(chart_shape.chart)
        return slide
    
    def _style_title_text(self, shape):
        """Apply corporate styling to title text - preserve template styling when possible"""
        try:
            # Check if we're using a template (if template styling exists, preserve it)
            has_template_styling = False
            
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.runs:
                    for run in paragraph.runs:
                        # Check if template has already set styling
                        if (run.font.name and run.font.name != "Calibri") or \
                           (run.font.size and run.font.size.pt > 10):
                            has_template_styling = True
                            break
                if has_template_styling:
                    break
            
            # Only apply default styling if no template styling is detected
            if not has_template_styling:
                print("No template styling detected, applying default styling")
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(32)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue
            else:
                print("Template styling detected, preserving template formatting")
                
        except Exception as e:
            print(f"Warning: Could not apply title styling: {e}")
            # Template styling will be preserved
    
    def _style_bullet_text(self, paragraph):
        """Apply corporate styling to bullet text - preserve template styling when possible"""
        try:
            # Check if template styling exists
            has_template_styling = False
            for run in paragraph.runs:
                if (run.font.name and run.font.name != "Calibri") or \
                   (run.font.size and run.font.size.pt > 10):
                    has_template_styling = True
                    break
            
            # Only apply default styling if no template styling is detected
            if not has_template_styling:
                print("No template bullet styling detected, applying default styling")
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(52, 73, 94)  # Medium blue
            else:
                print("Template bullet styling detected, preserving template formatting")
                
        except Exception as e:
            print(f"Warning: Could not apply bullet styling: {e}")
            # Template styling will be preserved
    
    def _style_table_header(self, cell):
        """Style table header cells"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(44, 62, 80)  # Dark blue
        
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)  # White
                run.font.bold = True
                run.font.size = Pt(14)
    
    def _style_table_cell(self, cell):
        """Style regular table cells"""
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(52, 73, 94)  # Medium blue
                run.font.size = Pt(12)
    
    def _style_chart(self, chart):
        """Apply corporate styling to charts"""
        # This would apply custom chart styling
        # For now, using default styling
        pass
    
    def edit_presentation(self, pptx_path: str, updates: Dict[str, Any]) -> str:
        """
        Edit an existing PowerPoint presentation using AI-generated instructions
        """
        print(f"🎯 Editing presentation: {pptx_path}")
        print(f"📝 Updates: {updates}")
        
        # Load existing presentation
        prs = Presentation(pptx_path)
        
        # Apply updates based on edit instructions
        if "edits" in updates:
            self._apply_edit_instructions(prs, updates["edits"])
        else:
            # Legacy format support
            self._apply_legacy_updates(prs, updates)
        
        # Save edited presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"edited_presentation_{timestamp}.pptx"
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        
        print(f"✅ Saved edited presentation: {output_path}")
        return output_path
    
    def _apply_edit_instructions(self, prs: Presentation, edit_instructions: List[Dict[str, Any]]):
        """
        Apply AI-generated edit instructions to presentation
        """
        for edit in edit_instructions:
            try:
                slide_index = edit.get("slide_index", 1) - 1  # Convert to 0-based index
                
                if slide_index < 0 or slide_index >= len(prs.slides):
                    print(f"⚠️ Invalid slide index: {slide_index + 1}")
                    continue
                
                slide = prs.slides[slide_index]
                action = edit.get("action", "")
                target_element = edit.get("target_element", "")
                changes = edit.get("changes", {})
                
                print(f"🔧 Applying {action} to {target_element} on slide {slide_index + 1}")
                
                if action == "modify_content":
                    self._modify_slide_content(slide, target_element, changes)
                elif action == "add_content":
                    self._add_slide_content(slide, target_element, changes)
                elif action == "replace_content":
                    self._replace_slide_content(slide, target_element, changes)
                elif action == "delete_content":
                    self._delete_slide_content(slide, target_element, changes)
                elif action == "change_layout":
                    self._change_slide_layout(slide, changes)
                
            except Exception as e:
                print(f"❌ Error applying edit: {e}")
                continue
    
    def _modify_slide_content(self, slide, target_element: str, changes: Dict[str, Any]):
        """Modify existing content on a slide"""
        if target_element == "title":
            # Find and modify title
            title_found = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if this looks like a title (short text, likely in title placeholder)
                    if len(shape.text.strip()) < 200 and not "\n" in shape.text:
                        if "find" in changes and "replace" in changes:
                            if changes["find"].lower() in shape.text.lower():
                                shape.text = shape.text.replace(changes["find"], changes["replace"])
                                print(f"   ✅ Modified title: '{changes['find']}' → '{changes['replace']}'")
                                title_found = True
                                break
                        elif "new_content" in changes:
                            shape.text = changes["new_content"]
                            print(f"   ✅ Set title to: '{changes['new_content']}'")
                            title_found = True
                            break
            
            if not title_found:
                print(f"   ⚠️ Title not found for modification")
        
        elif target_element == "text":
            # Modify any text content with find/replace
            text_modified = False
            find_text = changes.get("find", "")
            replace_text = changes.get("replace", "")
            
            print(f"   🔍 Looking for text: '{find_text}' to replace with: '{replace_text}'")
            
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    original_text = shape.text
                    print(f"   📄 Shape {shape_idx + 1} text: '{original_text[:100]}{'...' if len(original_text) > 100 else ''}'")
                    
                    if "find" in changes and "replace" in changes:
                        find_text = changes["find"]
                        replace_text = changes["replace"]
                        
                        # Try exact match first
                        if find_text in original_text:
                            new_text = original_text.replace(find_text, replace_text)
                            shape.text = new_text
                            print(f"   ✅ Exact match - Modified text: '{find_text}' → '{replace_text}'")
                            text_modified = True
                            break
                        
                        # Try case-insensitive search
                        elif find_text.lower() in original_text.lower():
                            # Find the actual case in the original text
                            start_idx = original_text.lower().find(find_text.lower())
                            if start_idx != -1:
                                # Replace while preserving surrounding text
                                actual_text = original_text[start_idx:start_idx + len(find_text)]
                                new_text = original_text.replace(actual_text, replace_text)
                                shape.text = new_text
                                print(f"   ✅ Case-insensitive match - Modified text: '{actual_text}' → '{replace_text}'")
                                text_modified = True
                                break
                        
                        # Try partial word match
                        else:
                            words_to_find = find_text.lower().split()
                            text_words = original_text.lower().split()
                            
                            # Check if all words from find_text exist in the shape text
                            if all(word in text_words for word in words_to_find):
                                # Try to find a reasonable replacement
                                for word in words_to_find:
                                    if word in original_text.lower():
                                        # Find the actual case version
                                        word_start = original_text.lower().find(word)
                                        actual_word = original_text[word_start:word_start + len(word)]
                                        if word == words_to_find[-1]:  # Replace the last/main word
                                            new_text = original_text.replace(actual_word, replace_text)
                                            shape.text = new_text
                                            print(f"   ✅ Partial match - Modified text: '{actual_word}' → '{replace_text}'")
                                            text_modified = True
                                            break
                            
                            if text_modified:
                                break
                    
                    elif "new_content" in changes:
                        shape.text = changes["new_content"]
                        print(f"   ✅ Set text to: '{changes['new_content']}'")
                        text_modified = True
                        break
            
            if not text_modified:
                print(f"   ⚠️ Text '{find_text}' not found for modification")
                print(f"   💡 Available text in slide:")
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        print(f"      Shape {shape_idx + 1}: '{shape.text[:50]}{'...' if len(shape.text) > 50 else ''}'")
        
        elif target_element == "bullets":
            # Find text box with bullets and modify
            for shape in slide.shapes:
                if hasattr(shape, "text") and ("•" in shape.text or len(shape.text.split('\n')) > 1):
                    if "find" in changes and "replace" in changes:
                        if changes["find"].lower() in shape.text.lower():
                            shape.text = shape.text.replace(changes["find"], changes["replace"])
                            print(f"   ✅ Modified bullets: '{changes['find']}' → '{changes['replace']}'")
                            break
                    elif "new_content" in changes:
                        shape.text = changes["new_content"]
                        print(f"   ✅ Set bullets to: '{changes['new_content']}'")
                        break
                    elif "add_bullet" in changes:
                        shape.text += f"\n• {changes['add_bullet']}"
                        print(f"   ✅ Added bullet: '{changes['add_bullet']}'")
                        break
    
    def _add_slide_content(self, slide, target_element: str, changes: Dict[str, Any]):
        """Add new content to a slide"""
        if target_element == "chart":
            self._add_chart_to_slide(slide, changes)
        elif target_element == "table":
            self._add_table_to_slide(slide, changes)
        elif target_element == "bullets":
            self._add_bullets_to_slide(slide, changes)
        elif target_element == "text":
            self._add_text_to_slide(slide, changes)
    
    def _add_chart_to_slide(self, slide, changes: Dict[str, Any]):
        """Add a chart to the slide"""
        try:
            chart_type = changes.get("chart_type", "pie")
            chart_data = changes.get("chart_data", {})
            
            if not chart_data:
                # Default chart data
                chart_data = {
                    "categories": ["Category A", "Category B", "Category C"],
                    "values": [30, 45, 25]
                }
            
            # Position for chart (right side of slide)
            left = Inches(5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(3)
            
            # Create chart data
            chart_data_obj = CategoryChartData()
            categories = chart_data.get("categories", ["A", "B", "C"])
            values = chart_data.get("values", [30, 45, 25])
            
            chart_data_obj.categories = categories
            chart_data_obj.add_series('Series 1', values)
            
            # Add chart to slide
            if chart_type.lower() == "pie":
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.PIE, left, top, width, height, chart_data_obj
                ).chart
            elif chart_type.lower() == "bar":
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data_obj
                ).chart
            else:
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.PIE, left, top, width, height, chart_data_obj
                ).chart
            
            print(f"✅ Added {chart_type} chart to slide")
            
        except Exception as e:
            print(f"❌ Error adding chart: {e}")
    
    def _add_table_to_slide(self, slide, changes: Dict[str, Any]):
        """Add a table to the slide"""
        try:
            table_data = changes.get("table_data", {})
            headers = table_data.get("headers", ["Column 1", "Column 2"])
            rows = table_data.get("rows", [["Data 1", "Data 2"], ["Data 3", "Data 4"]])
            
            # Position for table
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(3)
            
            # Create table
            rows_count = len(rows) + 1  # +1 for header
            cols_count = len(headers)
            
            table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
            table = table_shape.table
            
            # Add headers
            for i, header in enumerate(headers):
                table.cell(0, i).text = header
            
            # Add data rows
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < cols_count:
                        table.cell(row_idx + 1, col_idx).text = str(cell_data)
            
            print(f"✅ Added table with {rows_count} rows and {cols_count} columns")
            
        except Exception as e:
            print(f"❌ Error adding table: {e}")
    
    def _add_bullets_to_slide(self, slide, changes: Dict[str, Any]):
        """Add bullet points to the slide"""
        try:
            bullets = changes.get("bullets", [])
            new_content = changes.get("new_content", "")
            
            if new_content:
                bullets = [new_content]
            
            # Find existing text box or create new one
            text_shape = None
            for shape in slide.shapes:
                if hasattr(shape, "text") and ("•" in shape.text or len(shape.text.split('\n')) > 1):
                    text_shape = shape
                    break
            
            if text_shape:
                # Add to existing bullets
                for bullet in bullets:
                    text_shape.text += f"\n• {bullet}"
            else:
                # Create new text box
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4)
                
                text_shape = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_shape.text_frame
                text_frame.clear()
                
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = f"• {bullet}"
            
            print(f"✅ Added {len(bullets)} bullet points")
            
        except Exception as e:
            print(f"❌ Error adding bullets: {e}")
    
    def _add_text_to_slide(self, slide, changes: Dict[str, Any]):
        """Add text content to the slide"""
        try:
            new_text = changes.get("new_content", "New text content")
            
            # Position for new text
            left = Inches(1)
            top = Inches(4)
            width = Inches(8)
            height = Inches(2)
            
            text_shape = slide.shapes.add_textbox(left, top, width, height)
            text_shape.text = new_text
            
            print(f"✅ Added text content")
            
        except Exception as e:
            print(f"❌ Error adding text: {e}")
    
    def _replace_slide_content(self, slide, target_element: str, changes: Dict[str, Any]):
        """Replace existing content on a slide"""
        # First delete the target element
        self._delete_slide_content(slide, target_element, changes)
        # Then add new content
        self._add_slide_content(slide, target_element, changes)
    
    def _delete_slide_content(self, slide, target_element: str, changes: Dict[str, Any]):
        """Delete content from a slide"""
        shapes_to_remove = []
        
        for shape in slide.shapes:
            should_remove = False
            
            if target_element == "chart" and hasattr(shape, 'chart'):
                should_remove = True
            elif target_element == "table" and hasattr(shape, 'table'):
                should_remove = True
            elif target_element == "text" and hasattr(shape, 'text'):
                if "containing" in changes:
                    if changes["containing"].lower() in shape.text.lower():
                        should_remove = True
                else:
                    should_remove = True
            
            if should_remove:
                shapes_to_remove.append(shape)
        
        # Remove shapes (in reverse order to maintain indices)
        for shape in reversed(shapes_to_remove):
            slide.shapes.element.remove(shape.element)
        
        print(f"✅ Removed {len(shapes_to_remove)} {target_element} elements")
    
    def _change_slide_layout(self, slide, changes: Dict[str, Any]):
        """Change the layout of a slide"""
        # This is complex and would require layout matching
        # For now, just log the request
        layout_name = changes.get("new_layout", "unknown")
        print(f"📝 Layout change requested to: {layout_name}")
    
    def _apply_legacy_updates(self, prs: Presentation, updates: Dict[str, Any]):
        """Apply legacy update format for backwards compatibility"""
        # This handles the old update format
        for slide_idx, slide_updates in updates.items():
            if slide_idx.isdigit():
                slide_index = int(slide_idx) - 1
                if 0 <= slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    
                    if "title" in slide_updates:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip() and len(shape.text) < 100:
                                shape.text = slide_updates["title"]
                                break 

    def create_presentation(self, slide_data: Dict[str, Any], output_dir: Optional[str] = None, 
                           template_path: Optional[str] = None, logo_path: Optional[str] = None,
                           logo_position: str = "top-right", logo_size: str = "medium") -> str:
        """
        Legacy method - redirects to comprehensive template method
        """
        return self.create_presentation_with_full_template(
            slide_data, template_path, logo_path, logo_position, logo_size, output_dir
        )

    def extract_slide_data_from_ppt(self, ppt_path: str) -> Dict[str, Any]:
        """
        Extract slide data from an existing PowerPoint file for editing
        """
        print(f"🔍 Extracting slide data from: {ppt_path}")
        
        try:
            prs = Presentation(ppt_path)
            
            # Extract presentation metadata and SAVE TEMPLATE INFO
            slide_data = {
                "meta": {
                    "deck_title": "Extracted Presentation",
                    "original_file": os.path.basename(ppt_path),
                    "total_slides": len(prs.slides),
                    "extracted_on": datetime.now().isoformat(),
                    "original_template_path": ppt_path,  # Store original file path
                    "has_template": True  # Flag to indicate this needs template preservation
                },
                "slides": []
            }
            
            print(f"📊 Found {len(prs.slides)} slides to extract")
            
            for slide_idx, slide in enumerate(prs.slides):
                print(f"📄 Extracting slide {slide_idx + 1}...")
                
                slide_info = {
                    "title": "Untitled Slide",
                    "layout": "bullets",  # Default layout
                    "bullets": [],
                    "subtitle": "",
                    "columns": [],
                    "rows": [],
                    "categories": [],
                    "series": {},
                    "text_boxes": [],  # All other text content
                    "images": [],      # Image information
                    "shapes": []       # Other shapes
                }
                
                # Extract title
                try:
                    if hasattr(slide.shapes, 'title') and slide.shapes.title:
                        slide_info["title"] = slide.shapes.title.text or f"Slide {slide_idx + 1}"
                        print(f"   📋 Title: {slide_info['title']}")
                except:
                    slide_info["title"] = f"Slide {slide_idx + 1}"
                
                # Analyze slide content
                has_table = False
                has_chart = False
                text_content = []
                shape_count = 0
                
                print(f"   🔍 Analyzing {len(slide.shapes)} shapes on slide...")
                
                for shape in slide.shapes:
                    shape_count += 1
                    try:
                        print(f"      Shape {shape_count}: Type={shape.shape_type}")
                        
                        # Check for tables
                        if hasattr(shape, 'table'):
                            has_table = True
                            table = shape.table
                            num_rows = len(table.rows)
                            num_cols = len(table.columns)
                            print(f"   📊 Found table with {num_rows} rows, {num_cols} columns")
                            
                            # Extract table data
                            slide_info["layout"] = "table"
                            
                            # Extract headers (first row)
                            if num_rows > 0:
                                slide_info["columns"] = []
                                for col_idx in range(num_cols):
                                    try:
                                        header_text = table.cell(0, col_idx).text.strip()
                                        slide_info["columns"].append(header_text or f"Column {col_idx + 1}")
                                    except:
                                        slide_info["columns"].append(f"Column {col_idx + 1}")
                            
                            # Extract data rows
                            slide_info["rows"] = []
                            for row_idx in range(1, num_rows):  # Skip header row
                                row_data = []
                                for col_idx in range(num_cols):
                                    try:
                                        cell_text = table.cell(row_idx, col_idx).text.strip()
                                        row_data.append(cell_text)
                                    except:
                                        row_data.append("")
                                slide_info["rows"].append(row_data)
                        
                        # Check for charts
                        elif hasattr(shape, 'chart'):
                            has_chart = True
                            print(f"   📈 Found chart")
                            slide_info["layout"] = "chart.column"  # Default chart type
                            
                            # Try to extract chart data (basic extraction)
                            slide_info["categories"] = ["Category 1", "Category 2", "Category 3"]
                            slide_info["series"] = {"Series 1": [10, 20, 30]}
                        
                        # Check for images
                        elif hasattr(shape, 'image'):
                            print(f"   🖼️ Found image")
                            slide_info["images"].append({
                                "name": f"Image {len(slide_info['images']) + 1}",
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height
                            })
                        
                        # Extract ALL text content from ANY shape with text
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                print(f"   📝 Text found: '{text[:100]}{'...' if len(text) > 100 else ''}'")
                                
                                # Check if this is the title shape
                                try:
                                    if hasattr(slide.shapes, 'title') and slide.shapes.title == shape:
                                        slide_info["title"] = text
                                        print(f"   📋 Identified as title: {text}")
                                        continue
                                except:
                                    pass
                                
                                # Check if it's a placeholder and what type
                                is_placeholder = False
                                placeholder_type = None
                                try:
                                    if hasattr(shape, 'placeholder_format'):
                                        is_placeholder = True
                                        placeholder_type = shape.placeholder_format.type
                                        print(f"   🏷️ Placeholder type: {placeholder_type}")
                                        
                                        if placeholder_type == 1:  # Title placeholder
                                            slide_info["title"] = text
                                            print(f"   📋 Title placeholder: {text}")
                                            continue
                                        elif placeholder_type == 3:  # Subtitle placeholder
                                            slide_info["subtitle"] = text
                                            print(f"   📝 Subtitle placeholder: {text}")
                                            continue
                                        elif placeholder_type == 2:  # Content placeholder
                                            # This might be bullet points or other content
                                            if '\n' in text:
                                                bullets = [line.strip() for line in text.split('\n') if line.strip()]
                                                text_content.extend(bullets)
                                                print(f"   🔸 Content placeholder (bullets): {len(bullets)} items")
                                            else:
                                                text_content.append(text)
                                                print(f"   🔸 Content placeholder (single): {text}")
                                            continue
                                except:
                                    pass
                                
                                # For non-placeholder text or unknown placeholders, store as text box
                                text_box_info = {
                                    "text": text,
                                    "left": shape.left,
                                    "top": shape.top,
                                    "width": shape.width,
                                    "height": shape.height,
                                    "is_placeholder": is_placeholder,
                                    "placeholder_type": placeholder_type,
                                    "shape_id": f"textbox_{len(slide_info['text_boxes'])}"
                                }
                                
                                # Try to extract formatting information
                                try:
                                    if shape.text_frame.paragraphs:
                                        first_para = shape.text_frame.paragraphs[0]
                                        if first_para.runs:
                                            first_run = first_para.runs[0]
                                            text_box_info["font_name"] = getattr(first_run.font, 'name', None)
                                            text_box_info["font_size"] = getattr(first_run.font.size, 'pt', None) if first_run.font.size else None
                                            text_box_info["bold"] = getattr(first_run.font, 'bold', None)
                                            text_box_info["italic"] = getattr(first_run.font, 'italic', None)
                                except:
                                    pass
                                
                                slide_info["text_boxes"].append(text_box_info)
                                print(f"   📄 Added text box: '{text[:50]}{'...' if len(text) > 50 else ''}'")
                        
                        # Store other shape information
                        else:
                            shape_info = {
                                "type": str(shape.shape_type),
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height,
                                "shape_id": f"shape_{len(slide_info['shapes'])}"
                            }
                            slide_info["shapes"].append(shape_info)
                            print(f"   🔷 Added shape: {shape.shape_type}")
                    
                    except Exception as e:
                        print(f"   ❌ Error processing shape {shape_count}: {e}")
                        import traceback
                        traceback.print_exc()
                        continue
                
                # Set layout based on content found
                if has_table:
                    slide_info["layout"] = "table"
                elif has_chart:
                    slide_info["layout"] = "chart.column"
                elif slide_idx == 0:  # First slide is often a title slide
                    slide_info["layout"] = "title"
                else:
                    slide_info["layout"] = "bullets"
                    slide_info["bullets"] = text_content if text_content else ["Bullet point content"]
                
                print(f"   ✅ Layout determined: {slide_info['layout']}")
                slide_data["slides"].append(slide_info)
            
            print(f"✅ Successfully extracted {len(slide_data['slides'])} slides")
            return slide_data
            
        except Exception as e:
            print(f"❌ Error extracting slide data: {e}")
            import traceback
            traceback.print_exc()
            raise Exception(f"Failed to extract slide data: {str(e)}")

    def _update_existing_slide_content(self, slide, slide_info):
        """
        Update existing slide content while preserving template formatting
        """
        print(f"🔄 Updating existing slide content...")
        
        try:
            # Update title if it exists
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide.shapes.title.text = slide_info.get("title", "")
                print(f"   ✅ Updated title: {slide_info.get('title', '')}")
            
            # Find and update content placeholders and text boxes
            content_updated = False
            
            # Update subtitle if exists
            if slide_info.get("subtitle"):
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type == 3:  # Subtitle placeholder
                            shape.text = slide_info["subtitle"]
                            print(f"   ✅ Updated subtitle: {slide_info['subtitle']}")
                            break
            
            # Handle bullet points
            if slide_info.get("bullets"):
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type == 2:  # Content placeholder
                            # Update bullet content while preserving formatting
                            text_frame = shape.text_frame
                            text_frame.clear()
                            
                            for i, bullet in enumerate(slide_info["bullets"]):
                                if i == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                p.text = bullet
                                p.level = 0
                            
                            content_updated = True
                            print(f"   ✅ Updated {len(slide_info['bullets'])} bullet points")
                            break
            
            # Handle tables
            if slide_info.get("columns") and slide_info.get("rows"):
                # Remove existing table if any
                shapes_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, 'table'):
                        shapes_to_remove.append(shape)
                
                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)
                
                # Create new table
                rows = len(slide_info["rows"]) + 1  # +1 for header
                cols = len(slide_info["columns"])
                
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4)
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Set headers
                for i, header in enumerate(slide_info["columns"]):
                    table.cell(0, i).text = header
                
                # Set data rows
                for row_idx, row_data in enumerate(slide_info["rows"]):
                    for col_idx, cell_data in enumerate(row_data):
                        table.cell(row_idx + 1, col_idx).text = str(cell_data)
                
                content_updated = True
                print(f"   ✅ Updated table with {rows} rows and {cols} columns")
            
            # Handle text boxes
            if slide_info.get("text_boxes"):
                # Remove existing non-placeholder text shapes
                shapes_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # Keep placeholders, remove others
                        if not hasattr(shape, 'placeholder_format'):
                            shapes_to_remove.append(shape)
                
                for shape in shapes_to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)
                
                # Add updated text boxes
                self._add_text_boxes_to_slide(slide, slide_info)
                print(f"   ✅ Updated {len(slide_info['text_boxes'])} text boxes")
            
            if not content_updated:
                print(f"   ⚠️ No content placeholders found for update")
                
        except Exception as e:
            print(f"   ❌ Error updating slide content: {e}")
            import traceback
            traceback.print_exc()

    def _create_title_slide_with_template(self, prs, slide_info, template_assets):
        """Create title slide with full template inheritance"""
        slide_layout = self._get_slide_layout(prs, 'title')
        slide = prs.slides.add_slide(slide_layout)
        
        print(f"📋 Title slide using layout: {slide_layout.name}")
        
        # Set title using template placeholders
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
            print("✅ Used slide.shapes.title")
        else:
            # Find title placeholder
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:
                    placeholder.text = slide_info["title"]
                    print("✅ Used title placeholder (type 1)")
                    break
        
        # Set subtitle if present
        if "subtitle" in slide_info and slide_info["subtitle"]:
            subtitle_set = False
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 3:
                    placeholder.text = slide_info["subtitle"]
                    print("✅ Used subtitle placeholder (type 3)")
                    subtitle_set = True
                    break
            
            if not subtitle_set and len(slide.placeholders) > 1:
                try:
                    slide.placeholders[1].text = slide_info["subtitle"]
                    print("✅ Used second placeholder for subtitle")
                except:
                    pass
        
        return slide

    def _create_bullet_slide_with_template(self, prs, slide_info, template_assets):
        """Create bullet slide with full template inheritance"""
        slide_layout = self._get_slide_layout(prs, 'content')
        slide = prs.slides.add_slide(slide_layout)
        
        print(f"📋 Bullet slide using layout: {slide_layout.name}")
        
        # Set title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
            print("✅ Used slide.shapes.title for bullets")
        
        # Set bullet points using template placeholders
        content_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 2:  # Content placeholder
                content_placeholder = placeholder
                print("✅ Found content placeholder (type 2)")
                break
        
        if not content_placeholder and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            print("✅ Using second placeholder for content")
        
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(slide_info["bullets"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                # Don't override template styling - let it inherit
        
        return slide

    def _create_table_slide_with_template(self, prs, slide_info, template_assets):
        """Create table slide with template inheritance"""
        slide_layout = self._get_slide_layout(prs, 'content')
        slide = prs.slides.add_slide(slide_layout)
        
        print(f"📊 Table slide using layout: {slide_layout.name}")
        
        # Set title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
            print("✅ Used slide.shapes.title for table")
        
        # Remove content placeholder to avoid "Click to add text"
        placeholders_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 2:  # Content placeholder
                    placeholders_to_remove.append(shape)
                    print("🗑️ Removing content placeholder to clean slide")
        
        # Remove content placeholders
        for placeholder in placeholders_to_remove:
            try:
                sp = placeholder._element
                sp.getparent().remove(sp)
            except:
                pass  # Continue if removal fails
        
        # Create table
        rows = len(slide_info["rows"]) + 1  # +1 for header
        cols = len(slide_info["columns"])
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set headers with template styling
        for i, header in enumerate(slide_info["columns"]):
            cell = table.cell(0, i)
            cell.text = header
            # Apply header styling that matches template
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(70, 130, 180)  # Steel blue to match template
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                    run.font.bold = True
        
        # Set data rows
        for row_idx, row_data in enumerate(slide_info["rows"]):
            for col_idx, cell_data in enumerate(row_data):
                table.cell(row_idx + 1, col_idx).text = str(cell_data)
        
        return slide

    def _create_chart_slide_with_template(self, prs, slide_info, template_assets):
        """Create chart slide with template inheritance"""
        slide_layout = self._get_slide_layout(prs, 'content')
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = slide_info["title"]
        
        # Add chart placeholder text for now
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        text_frame = text_box.text_frame
        text_frame.text = f"Chart: {slide_info['title']}\nCategories: {', '.join(slide_info.get('categories', []))}\nData available for chart generation"
        
        return slide

    def _apply_template_assets_to_slide(self, slide, template_assets):
        """
        Apply extracted template assets (logos, etc.) to a slide
        """
        if not template_assets or 'logos' not in template_assets:
            return
        
        print(f"🖼️  Applying {len(template_assets['logos'])} template assets to slide...")
        
        for logo_info in template_assets['logos']:
            try:
                if os.path.exists(logo_info['path']):
                    # Add the extracted logo/image to the slide at its original position
                    slide.shapes.add_picture(
                        logo_info['path'],
                        logo_info['left'],
                        logo_info['top'],
                        logo_info['width'],
                        logo_info['height']
                    )
                    print(f"   ✅ Added {logo_info['filename']} to slide")
            except Exception as e:
                print(f"   ❌ Could not add {logo_info['filename']}: {e}")

    def _add_text_boxes_to_slide(self, slide, slide_info):
        """
        Add text boxes from extracted slide data
        """
        if 'text_boxes' not in slide_info or not slide_info['text_boxes']:
            return
        
        print(f"📄 Adding {len(slide_info['text_boxes'])} text boxes to slide...")
        
        for i, text_box in enumerate(slide_info['text_boxes']):
            try:
                # Get position and size, with defaults
                left = text_box.get('left', Inches(1))
                top = text_box.get('top', Inches(2))
                width = text_box.get('width', Inches(6))
                height = text_box.get('height', Inches(1))
                text = text_box.get('text', '')
                
                if not text.strip():
                    continue
                
                # Create text box
                text_box_shape = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box_shape.text_frame
                text_frame.text = text
                
                # Apply formatting if available
                try:
                    if text_frame.paragraphs:
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                # Apply font formatting from extracted data
                                if text_box.get('font_name'):
                                    run.font.name = text_box['font_name']
                                if text_box.get('font_size'):
                                    run.font.size = Pt(text_box['font_size'])
                                if text_box.get('bold'):
                                    run.font.bold = text_box['bold']
                                if text_box.get('italic'):
                                    run.font.italic = text_box['italic']
                except Exception as format_error:
                    print(f"   ⚠️ Could not apply formatting to text box {i+1}: {format_error}")
                
                print(f"   ✅ Added text box {i+1}: '{text[:50]}{'...' if len(text) > 50 else ''}'")
                
            except Exception as e:
                print(f"   ❌ Could not add text box {i+1}: {e}")
                import traceback
                traceback.print_exc()

    def extract_slide_content(self, pptx_path):
        """
        Extract content from all slides for AI analysis
        """
        try:
            prs = Presentation(pptx_path)
            slide_data = {
                "meta": {
                    "total_slides": len(prs.slides),
                    "title": "Extracted Presentation"
                },
                "slides": []
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_number": slide_idx + 1,
                    "layout": f"layout_{slide.slide_layout.slide_layout.slide_layout.slide_index if hasattr(slide.slide_layout, 'slide_layout') else 0}",
                    "content": []
                }
                
                # Extract content from each shape
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_info = {
                        "shape_id": shape_idx,
                        "type": str(shape.shape_type)
                    }
                    
                    # Extract text content
                    if hasattr(shape, 'text') and shape.text:
                        shape_info["text"] = shape.text
                        shape_info["type"] = "text"
                        
                        # Try to determine if it's a title or content
                        if shape_idx == 0 or 'title' in str(shape.name).lower():
                            shape_info["is_title"] = True
                        else:
                            shape_info["is_title"] = False
                    
                    # Extract table content
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        shape_info["type"] = "table"
                        try:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = []
                                for cell in row.cells:
                                    row_data.append(cell.text.strip())
                                table_data.append(row_data)
                            shape_info["table_data"] = table_data
                        except Exception as e:
                            shape_info["table_data"] = f"Table extraction error: {e}"
                    
                    # Extract chart info
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        shape_info["type"] = "chart"
                        shape_info["chart_type"] = "chart_element"
                    
                    # Add shape if it has content
                    if any(key in shape_info for key in ["text", "table_data", "chart_type"]):
                        slide_content["content"].append(shape_info)
                
                slide_data["slides"].append(slide_content)
            
            return slide_data
            
        except Exception as e:
            print(f"Error extracting slide content: {e}")
            import traceback
            traceback.print_exc()
            raise

    def preview_edits(self, pptx_path, edit_instructions, target_slide_number=None):
        """
        Generate a preview of edits without actually modifying the file
        Returns a preview object showing what will change
        """
        try:
            # Load presentation
            prs = Presentation(pptx_path)
            preview_data = {
                "slides": [],
                "edit_summary": [],
                "total_slides": len(prs.slides)
            }
            
            # Convert slide number to index (1-based to 0-based)
            target_slide_index = None
            if target_slide_number is not None:
                target_slide_index = target_slide_number - 1
                if target_slide_index < 0 or target_slide_index >= len(prs.slides):
                    raise ValueError(f"Slide {target_slide_number} does not exist (total slides: {len(prs.slides)})")
            
            # Process each slide or just the target slide
            slides_to_process = [target_slide_index] if target_slide_index is not None else range(len(prs.slides))
            
            for slide_idx in slides_to_process:
                slide = prs.slides[slide_idx]
                slide_preview = {
                    "slide_number": slide_idx + 1,
                    "original_content": [],
                    "changes": [],
                    "new_content": []
                }
                
                # Extract current content
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_info = {
                        "shape_id": shape_idx,
                        "type": str(shape.shape_type),
                        "content": ""
                    }
                    
                    if hasattr(shape, 'text') and shape.text:
                        shape_info["content"] = shape.text
                        shape_info["type"] = "text"
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        shape_info["type"] = "table"
                        shape_info["content"] = self._extract_table_preview(shape)
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        shape_info["type"] = "chart"
                        shape_info["content"] = "Chart element"
                    
                    slide_preview["original_content"].append(shape_info)
                
                # Process edit instructions for this slide
                for instruction in edit_instructions:
                    if self._should_apply_to_slide(instruction, slide_idx + 1, target_slide_number):
                        change_preview = self._preview_instruction(instruction, slide_preview["original_content"])
                        if change_preview:
                            slide_preview["changes"].append(change_preview)
                            preview_data["edit_summary"].append(f"Slide {slide_idx + 1}: {change_preview['description']}")
                
                preview_data["slides"].append(slide_preview)
            
            return preview_data
            
        except Exception as e:
            print(f"Error generating preview: {e}")
            import traceback
            traceback.print_exc()
            raise
    
    def _extract_table_preview(self, table_shape):
        """Extract table content for preview"""
        try:
            table = table_shape.table
            content = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                content.append(row_data)
            return content
        except Exception as e:
            return f"Table (preview error: {e})"
    
    def _preview_instruction(self, instruction, original_content):
        """Generate a preview of what an instruction will change"""
        try:
            action = instruction.get('action', '').lower()
            
            if action in ['modify_text', 'replace_text']:
                old_text = instruction.get('old_text', '')
                new_text = instruction.get('new_text', '')
                
                # Find matching content
                for content_item in original_content:
                    if old_text.lower() in content_item.get('content', '').lower():
                        return {
                            "action": "text_change",
                            "description": f"Change '{old_text}' to '{new_text}'",
                            "old_value": old_text,
                            "new_value": new_text,
                            "shape_id": content_item['shape_id']
                        }
            
            elif action == 'add_text':
                return {
                    "action": "add_content",
                    "description": f"Add text: '{instruction.get('text', '')}'",
                    "new_value": instruction.get('text', ''),
                    "type": "text"
                }
            
            elif action in ['add_chart', 'add_pie_chart']:
                return {
                    "action": "add_content", 
                    "description": f"Add {instruction.get('chart_type', 'chart')}",
                    "new_value": f"{instruction.get('chart_type', 'Chart')} with data",
                    "type": "chart"
                }
            
            elif action == 'add_table':
                return {
                    "action": "add_content",
                    "description": "Add table",
                    "new_value": "New table",
                    "type": "table"
                }
            
            return {
                "action": "unknown",
                "description": f"Apply {action}",
                "details": instruction
            }
            
        except Exception as e:
            return {
                "action": "error",
                "description": f"Preview error: {e}",
                "details": instruction
            }
    
    def _should_apply_to_slide(self, instruction, current_slide_number, target_slide_number):
        """
        Determine if an instruction should be applied to a specific slide
        """
        try:
            # Handle case where instruction might be a string instead of dict
            if isinstance(instruction, str):
                return True  # Apply to current slide if instruction is just a string
                
            # If a target slide is specified, only apply to that slide
            if target_slide_number is not None:
                return current_slide_number == target_slide_number
            
            # Check if instruction specifies a slide number
            slide_number = instruction.get('slide_number') if hasattr(instruction, 'get') else None
            if slide_number is not None:
                return current_slide_number == slide_number
            
            # If no specific slide mentioned, apply to current slide
            return True
            
        except Exception as e:
            print(f"Error checking slide application: {e}")
            return True  # Default to applying the instruction 